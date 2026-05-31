"""
演示文稿处理器

处理 PPTX 等演示文稿格式
"""

from pathlib import Path
from typing import Any, Dict, List
from .base_processor import BaseFileProcessor


class PresentationProcessor(BaseFileProcessor):
    """PPTX 演示文稿处理器"""
    
    async def to_text(self, file_path: Path) -> str:
        """将演示文稿转换为文本表示"""
        
        try:
            from pptx import Presentation
        except ImportError:
            return "[错误] 需要安装 python-pptx: pip install python-pptx"
        
        prs = Presentation(file_path)
        
        text_parts = [
            f"[演示文稿: {file_path.name}]\n\n",
            f"幻灯片数量: {len(prs.slides)}\n\n",
        ]
        
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n{'='*50}\n")
            text_parts.append(f"幻灯片 {i}\n")
            text_parts.append(f"{'='*50}\n")
            
            # 提取所有文本
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            # 检测是否为标题
                            if para.level == 0 and shape == slide.shapes[0]:
                                text_parts.append(f"\n## {text}\n")
                            else:
                                indent = "  " * para.level
                                bullet = "• " if para.level > 0 else ""
                                text_parts.append(f"{indent}{bullet}{text}\n")
                
                # 处理表格
                if shape.has_table:
                    table = shape.table
                    text_parts.append("\n[表格]\n")
                    for row in table.rows:
                        cells = [cell.text for cell in row.cells]
                        text_parts.append("| " + " | ".join(cells) + " |\n")
            
            # 添加分隔
            text_parts.append("\n")
        
        return "".join(text_parts)
    
    async def from_text(self, text: str, output_path: Path) -> bool:
        """从文本描述生成 PPTX"""
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return False
        
        prs = Presentation()
        
        # 解析幻灯片结构
        slides_content = self._parse_slides(text)
        
        for slide_content in slides_content:
            # 使用标题和内容布局
            slide_layout = prs.slide_layouts[1]  # 标题和内容布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 设置标题
            if slide_content.get('title'):
                title = slide.shapes.title
                title.text = slide_content['title']
            
            # 设置内容
            if slide_content.get('content'):
                body = slide.placeholders[1]
                tf = body.text_frame
                
                for i, item in enumerate(slide_content['content']):
                    if i == 0:
                        tf.text = item
                    else:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 0
        
        # 如果没有解析到幻灯片，创建一个默认的
        if len(prs.slides) == 0:
            slide_layout = prs.slide_layouts[0]  # 标题幻灯片
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "演示文稿"
        
        prs.save(output_path)
        return True
    
    def _parse_slides(self, text: str) -> List[Dict[str, Any]]:
        """解析幻灯片文本描述"""
        
        slides = []
        current_slide = None
        
        for line in text.split('\n'):
            line = line.strip()
            
            # 检测幻灯片标题 (## 开头)
            if line.startswith('## '):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {'title': line[3:], 'content': []}
            
            # 检测内容项
            elif current_slide:
                if line.startswith('- ') or line.startswith('• '):
                    current_slide['content'].append(line[2:])
                elif line and not line.startswith('='):
                    current_slide['content'].append(line)
        
        # 添加最后一个幻灯片
        if current_slide:
            slides.append(current_slide)
        
        return slides
    
    def get_supported_extensions(self) -> list[str]:
        return ['.pptx', '.ppt']
