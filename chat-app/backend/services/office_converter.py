import json
import os
from typing import Dict, List, Any, Optional
from docx import Document as DocxDocument
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook as OpenpyxlWorkbook, load_workbook
from pptx import Presentation as PptxPresentation
from pptx.util import Inches as PptxInches


class OfficeConverter:
    """Office 文件与 JSON 双向转换器."""

    @staticmethod
    def docx_to_json(file_path: str) -> Dict[str, Any]:
        doc = DocxDocument(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            runs = []
            for run in para.runs:
                runs.append({
                    "text": run.text,
                    "bold": run.bold,
                    "italic": run.italic,
                    "underline": run.underline,
                    "font_size": run.font.size.pt if run.font.size else None,
                })
            style_name = para.style.name if para.style else "Normal"
            align = None
            if para.alignment == WD_ALIGN_PARAGRAPH.CENTER:
                align = "center"
            elif para.alignment == WD_ALIGN_PARAGRAPH.RIGHT:
                align = "right"
            elif para.alignment == WD_ALIGN_PARAGRAPH.JUSTIFY:
                align = "justify"
            paragraphs.append({
                "text": para.text,
                "runs": runs,
                "style": style_name,
                "alignment": align,
            })
        return {"type": "document", "paragraphs": paragraphs}

    @staticmethod
    def json_to_docx(data: Dict[str, Any], output_path: str) -> str:
        doc = DocxDocument()
        paragraphs = data.get("paragraphs", [])
        for para_data in paragraphs:
            p = doc.add_paragraph()
            align = para_data.get("alignment")
            if align == "center":
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif align == "right":
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            elif align == "justify":
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            for run_data in para_data.get("runs", []):
                run = p.add_run(run_data.get("text", ""))
                run.bold = run_data.get("bold", False)
                run.italic = run_data.get("italic", False)
                run.underline = run_data.get("underline", False)
                font_size = run_data.get("font_size")
                if font_size:
                    run.font.size = Pt(font_size)
        doc.save(output_path)
        return output_path

    @staticmethod
    def xlsx_to_json(file_path: str) -> Dict[str, Any]:
        wb = load_workbook(file_path, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            cells = {}
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        addr = cell.coordinate
                        cells[addr] = {
                            "value": str(cell.value) if cell.value is not None else "",
                            "formula": cell.value if isinstance(cell.value, str) and str(cell.value).startswith("=") else None,
                        }
            sheets.append({"name": sheet_name, "cells": cells})
        return {"type": "spreadsheet", "sheets": sheets, "activeSheetIndex": 0}

    @staticmethod
    def json_to_xlsx(data: Dict[str, Any], output_path: str) -> str:
        wb = OpenpyxlWorkbook()
        sheets_data = data.get("sheets", [])
        if not sheets_data:
            sheets_data = [{"name": "Sheet1", "cells": {}}]
        for idx, sheet_data in enumerate(sheets_data):
            if idx == 0:
                ws = wb.active
                ws.title = sheet_data.get("name", "Sheet1")
            else:
                ws = wb.create_sheet(title=sheet_data.get("name", f"Sheet{idx+1}"))
            for addr, cell_data in sheet_data.get("cells", {}).items():
                ws[addr] = cell_data.get("value", "")
        wb.save(output_path)
        return output_path

    @staticmethod
    def pptx_to_json(file_path: str) -> Dict[str, Any]:
        prs = PptxPresentation(file_path)
        slides = []
        for slide in prs.slides:
            elements = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    elements.append({
                        "type": "text",
                        "x": shape.left,
                        "y": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                        "content": shape.text_frame.text,
                    })
            slides.append({"id": str(len(slides)), "elements": elements, "background": None})
        return {"type": "presentation", "slides": slides, "activeSlideIndex": 0}

    @staticmethod
    def json_to_pptx(data: Dict[str, Any], output_path: str) -> str:
        prs = PptxPresentation()
        slides_data = data.get("slides", [])
        for slide_data in slides_data:
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
            slide = prs.slides.add_slide(blank_layout)
            for el in slide_data.get("elements", []):
                if el.get("type") == "text":
                    left = el.get("x", 0)
                    top = el.get("y", 0)
                    width = el.get("width", PptxInches(2))
                    height = el.get("height", PptxInches(1))
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    txBox.text_frame.text = el.get("content", "")
        prs.save(output_path)
        return output_path

    @classmethod
    def convert_to_json(cls, file_path: str, file_type: str) -> Dict[str, Any]:
        if file_type == "document":
            return cls.docx_to_json(file_path)
        elif file_type == "spreadsheet":
            return cls.xlsx_to_json(file_path)
        elif file_type == "presentation":
            return cls.pptx_to_json(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    @classmethod
    def convert_from_json(cls, data: Dict[str, Any], output_path: str, file_type: str) -> str:
        if file_type == "document":
            return cls.json_to_docx(data, output_path)
        elif file_type == "spreadsheet":
            return cls.json_to_xlsx(data, output_path)
        elif file_type == "presentation":
            return cls.json_to_pptx(data, output_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
